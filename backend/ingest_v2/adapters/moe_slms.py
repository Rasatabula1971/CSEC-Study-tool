# PHASE: build
"""
backend/ingest_v2/adapters/moe_slms.py
======================================
Adapter for T&T MoE SLMS lesson material (.docx / .pptx / .pptm / .pdf) under a
subject's ``Notes\\T&T MoE SLMS`` folder.

SCOPE: only files under ``Notes\\T&T MoE SLMS`` are claimed. The same vendor folder
under ``Practice Questions`` and ``SBA`` is deliberately NOT handled here -- those
files have an uncharacterised format and the ``S{N} Obj {range}`` filename pattern
may not apply; they fall through to other adapters / review.

Objective mapping comes from the filename, e.g.:
    "S1 Obj 1-7"  -> section 1, objectives 1..7
    "S2 Obj 1 2"  -> section 2, objectives [1, 2]
    "S2 Obj 4-7"  -> section 2, objectives 4..7
    "S2 Obj 10"   -> section 2, objectives [10]
A filename with no parseable pattern yields ONE review record (no keyword guessing).
Every chunk binds to every objective in the parsed range (the corpus's many-to-many
case). confidence = high on a clean parse, medium when the objective spec needed
lenient coercion.
"""

import re
from pathlib import Path
from typing import Iterable, Optional

from backend.ingest_v2.adapters.base import BaseAdapter, IngestRecord, sha256_file
from backend.ingest_v2.chunking import chunk_paragraphs
from backend.ingest_v2.manifest import SubjectManifest
from backend.ingest_v2.normalize import normalize_text
from backend.ingest_v2.objective_index import ObjectiveIndex

# "S<n> ... Obj <spec>", spec = digits with spaces / hyphens / commas.
_MOE_RE = re.compile(r"\bS\s*(\d+)\b.*?\bObj\s*([0-9][0-9\s,\-]*)", re.IGNORECASE)
# Clean objective specs: a single number, a space-separated list, or an a-b range.
_CLEAN_RANGE_RE = re.compile(r"^\d+\s*-\s*\d+$")
_CLEAN_LIST_RE = re.compile(r"^\d+(?:\s+\d+)*$")

# The low-text PDF-page OCR trigger is the shared ingest_v2 threshold
# (ocr_utils.OCR_TRIGGER_THRESHOLD), imported at the call site in _extract_pdf rather
# than redefined here -- same value generic_pdf's opt-in OCR path uses.

# pptx/pptm and docx/pdf handled separately.
_PPT_SUFFIXES = {".pptx", ".pptm"}
_SUPPORTED = {".docx", ".pptx", ".pptm", ".pdf"}


def _parse_obj_spec(spec: str) -> tuple[list[int], bool]:
    """(objectives, clean). clean=True for the documented forms; clean=False when we
    fell back to a lenient digit scan. Returns ([], False) if no digits at all."""
    spec = spec.strip()
    if _CLEAN_RANGE_RE.match(spec):
        a, b = (int(x) for x in re.split(r"\s*-\s*", spec))
        lo, hi = (a, b) if a <= b else (b, a)
        return list(range(lo, hi + 1)), True
    if _CLEAN_LIST_RE.match(spec):
        return [int(x) for x in spec.split()], True
    nums = [int(x) for x in re.findall(r"\d+", spec)]
    if not nums:
        return [], False
    # Lenient: a two-number hyphenated form -> range; otherwise the literal list.
    if "-" in spec and len(nums) == 2:
        lo, hi = sorted(nums)
        return list(range(lo, hi + 1)), False
    return nums, False


def parse_moe_filename(stem: str) -> Optional[tuple[int, list[int], str]]:
    """(section, objectives, confidence) or None if the filename has no pattern."""
    m = _MOE_RE.search(stem)
    if not m:
        return None
    section = int(m.group(1))
    objs, clean = _parse_obj_spec(m.group(2))
    if not objs:
        return None
    return section, objs, ("high" if clean else "medium")


class MoESLMSAdapter(BaseAdapter):
    source_family = "moe_slms"

    def matches(self, path: Path) -> bool:
        if path.suffix.lower() not in _SUPPORTED:
            return False
        parts = set(path.parts)
        # Scoped to Notes\T&T MoE SLMS only (NOT Practice Questions / SBA).
        return "T&T MoE SLMS" in parts and "Notes" in parts

    def extract(self, path: Path, manifest: SubjectManifest,
                objective_index: ObjectiveIndex) -> Iterable[IngestRecord]:
        chash = sha256_file(path)
        subject_id = manifest.subject_id

        def _review(reason: str) -> IngestRecord:
            return IngestRecord(
                objective_id="REVIEW", subject_id=subject_id,
                source_family=self.source_family, content_type="notes",
                source_file=str(path), content_hash=chash, page=None,
                chunk_text=None, mcq_payload=None,
                confidence="review", review_reason=reason,
            )

        parsed = parse_moe_filename(path.stem)
        if parsed is None:
            yield _review("moe_slms_filename_unparsable")
            return
        section, obj_nums, confidence = parsed

        # Resolve the parsed objective numbers; drop (to review) any not in syllabus.
        valid_ids: list[str] = []
        for n in obj_nums:
            oid = objective_index.build_objective_id(section, n)
            if oid in objective_index.all_objective_ids():
                valid_ids.append(oid)
            else:
                yield IngestRecord(
                    objective_id=oid, subject_id=subject_id,
                    source_family=self.source_family, content_type="notes",
                    source_file=str(path), content_hash=chash, page=None,
                    chunk_text=None, mcq_payload=None,
                    confidence="review", review_reason="objective_id_not_in_syllabus",
                )
        if not valid_ids:
            return

        for page, raw in self._extract_pages(path):
            body = normalize_text(raw)
            for chunk in chunk_paragraphs(body):
                for oid in valid_ids:
                    yield IngestRecord(
                        objective_id=oid, subject_id=subject_id,
                        source_family=self.source_family, content_type="notes",
                        source_file=str(path), content_hash=chash, page=page,
                        chunk_text=chunk, mcq_payload=None,
                        confidence=confidence, review_reason=None,
                    )

    # --- text extraction by file type ------------------------------------
    def _extract_pages(self, path: Path):
        suffix = path.suffix.lower()
        if suffix == ".docx":
            yield from self._extract_docx(path)
        elif suffix in _PPT_SUFFIXES:
            yield from self._extract_pptx(path)
        elif suffix == ".pdf":
            yield from self._extract_pdf(path)

    @staticmethod
    def _extract_docx(path: Path):
        import docx  # python-docx
        doc = docx.Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        yield 1, "\n\n".join(paras)

    @staticmethod
    def _extract_pptx(path: Path):
        from pptx import Presentation
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text)
            yield i, "\n\n".join(texts)

    def _extract_pdf(self, path: Path):
        # backend/ on path already (imported by the package); reuse v1's PyMuPDF path.
        import sys
        sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
        import ingest as v1
        import ocr_utils  # shared OCR trigger threshold (backend/ on path above)
        for page, text in v1.extract_pdf_pages(path):
            if len((text or "").strip()) < ocr_utils.OCR_TRIGGER_THRESHOLD:
                ocr = self._ocr_pdf_page(path, page)
                if ocr:
                    text = ocr
            yield page, text

    @staticmethod
    def _ocr_pdf_page(path: Path, page: int) -> str:
        """Best-effort Tesseract OCR of one PDF page via the SHARED ocr_utils.ocr_page
        helper -- which adds the decompression-bomb DPI guard this adapter previously
        lacked (oversized pages no longer trip Pillow's guard). Any failure (no
        Tesseract, no Pillow, render error) returns '' so extraction degrades
        gracefully."""
        try:
            import sys
            sys.path.insert(0, str(Path(__file__).resolve().parents[2]))  # backend/ for ocr_utils
            import fitz  # PyMuPDF
            import ocr_utils
            doc = fitz.open(str(path))
            try:
                text, _conf = ocr_utils.ocr_page(doc.load_page(page - 1))
                return text or ""
            finally:
                doc.close()
        except Exception:
            return ""
