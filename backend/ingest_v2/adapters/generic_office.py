# PHASE: build
"""
backend/ingest_v2/adapters/generic_office.py
=============================================
Catch-all adapter for Office documents (.docx / .pptx / .pptm) that no specialised
adapter claimed -- lecture notes, purpose-built "Bridge"/"Supplemental" lesson
notes, etc. The Office analogue of GenericPDFAdapter: registered AFTER
MoESLMSAdapter in dispatch order, so a .docx/.pptx under ``Notes\\T&T MoE SLMS`` is
claimed by MoESLMSAdapter first and only the rest fall through here.

Objective mapping, in priority order (parse_office_filename), grounded in the real
filenames observed across subjects:
  1. Objective-id form  -- "POB-1.2 Bridge Lesson - ...", "ISCI-10.1 Bridge ...",
     "MATH-1.11 Bridge ..."  (the id already encodes section.objective)
  2. Underscore form     -- "S08_Obj3_Caribbean Economies ...", "S06_Obj8_Nominal ..."
  3. MoE space form      -- "CSEC Economics S2 Obj 10 ..." / "S1 Obj 1-7 ..."
     (reuses MoESLMSAdapter.parse_moe_filename -- not reimplemented)
A clean filename match binds every chunk to that objective (file-level, confidence
"high"). When NO pattern matches, we fall back to per-chunk keyword matching against
the ObjectiveIndex (the same v1 matcher GenericPDFAdapter uses): a match is
confidence "medium", a miss is routed to review. Unparseable files are never
silently skipped (Rule 1).

NOTE on prefixes: the objective-id form is parsed PREFIX-AGNOSTICALLY -- only the
section.objective is taken, and the canonical id is rebuilt with the SUBJECT's own
prefix (ObjectiveIndex), then validated against the locked syllabus. This is
deliberate: some subjects' Bridge files use a prefix (e.g. ISCI / ENGA) that differs
from the framework's provisional SUBJECT_PREFIX (INTSCI / ENG), so requiring the
filename prefix to match would wrongly reject them. The syllabus-membership check is
the Rule-1 guard against a bad (section, objective).
"""

import re
from pathlib import Path
from typing import Iterable, Optional

from backend.ingest_v2.adapters.base import BaseAdapter, IngestRecord, sha256_file
from backend.ingest_v2.adapters.moe_slms import parse_moe_filename
from backend.ingest_v2.chunking import chunk_paragraphs
from backend.ingest_v2.manifest import SubjectManifest
from backend.ingest_v2.normalize import normalize_text
from backend.ingest_v2.objective_index import ObjectiveIndex

_SUPPORTED = {".docx", ".pptx", ".pptm"}
_PPT_SUFFIXES = {".pptx", ".pptm"}

# 1. Objective-id: "<PREFIX>-<section>.<objective>" anywhere in the stem.
_OBJ_ID_RE = re.compile(r"\b[A-Za-z]{2,6}-(\d+)\.(\d+)\b")
# 2. Underscore: "S<section>_Obj<spec>", spec a single number or an a-b range.
_UNDERSCORE_RE = re.compile(r"\bS(\d+)_Obj\s*(\d+(?:\s*-\s*\d+)?)", re.IGNORECASE)


def _parse_underscore_objs(spec: str) -> list[int]:
    """'3' -> [3]; '3-4' -> [3, 4]. (The real Bridge/Supplemental files are
    single-objective; the range form is handled for safety.)"""
    spec = spec.strip()
    if "-" in spec:
        a, b = (int(x) for x in re.split(r"\s*-\s*", spec))
        lo, hi = (a, b) if a <= b else (b, a)
        return list(range(lo, hi + 1))
    return [int(spec)]


def parse_office_filename(stem: str) -> Optional[tuple[int, list[int], str]]:
    """(section, objectives, confidence) or None. Tries the three real conventions in
    priority order (objective-id, underscore, MoE space). confidence is 'high' for a
    clean structural match; parse_moe_filename's own 'medium' (lenient) is preserved."""
    m = _OBJ_ID_RE.search(stem)
    if m:
        return int(m.group(1)), [int(m.group(2))], "high"
    m = _UNDERSCORE_RE.search(stem)
    if m:
        return int(m.group(1)), _parse_underscore_objs(m.group(2)), "high"
    moe = parse_moe_filename(stem)   # reuse, not reimplement (MoE space form)
    if moe:
        return moe
    return None


class GenericOfficeAdapter(BaseAdapter):
    source_family = "generic_office"

    def matches(self, path: Path) -> bool:
        # Office formats only. Registered AFTER MoESLMSAdapter, so a .docx/.pptx under
        # Notes\T&T MoE SLMS is claimed there first; this catches everything else.
        return path.suffix.lower() in _SUPPORTED

    def extract(self, path: Path, manifest: SubjectManifest,
                objective_index: ObjectiveIndex) -> Iterable[IngestRecord]:
        chash = sha256_file(path)
        subject_id = manifest.subject_id

        def _record(oid, chunk, confidence, review_reason=None) -> IngestRecord:
            return IngestRecord(
                objective_id=oid, subject_id=subject_id,
                source_family=self.source_family, content_type="notes",
                source_file=str(path), content_hash=chash, page=None,
                chunk_text=chunk, mcq_payload=None,
                confidence=confidence, review_reason=review_reason,
            )

        pages = list(self._extract_pages(path))

        parsed = parse_office_filename(path.stem)
        if parsed is not None:
            section, obj_nums, confidence = parsed
            valid_ids: list[str] = []
            for n in obj_nums:
                oid = objective_index.build_objective_id(section, n)
                if oid in objective_index.all_objective_ids():
                    valid_ids.append(oid)
                else:
                    # filename named an objective that is not in the locked syllabus
                    yield IngestRecord(
                        objective_id=oid, subject_id=subject_id,
                        source_family=self.source_family, content_type="notes",
                        source_file=str(path), content_hash=chash, page=None,
                        chunk_text=None, mcq_payload=None,
                        confidence="review", review_reason="objective_id_not_in_syllabus",
                    )
            if valid_ids:
                for page, raw in pages:
                    for chunk in chunk_paragraphs(normalize_text(raw)):
                        for oid in valid_ids:
                            yield _record(oid, chunk, confidence)
            return

        # No filename pattern -> per-chunk keyword fallback (same matcher as GenericPDF).
        for page, raw in pages:
            for chunk in chunk_paragraphs(normalize_text(raw)):
                oid, _score = objective_index.resolve_by_keyword(chunk)
                if oid is None:
                    yield _record("REVIEW", chunk, "review",
                                  review_reason="no_objective_match_via_keywords")
                else:
                    yield _record(oid, chunk, "medium")

    # --- text extraction by file type (mirrors MoESLMSAdapter's docx/pptx readers) ---
    def _extract_pages(self, path: Path):
        suffix = path.suffix.lower()
        if suffix == ".docx":
            yield from self._extract_docx(path)
        elif suffix in _PPT_SUFFIXES:
            yield from self._extract_pptx(path)

    @staticmethod
    def _extract_docx(path: Path):
        import docx  # python-docx
        doc = docx.Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        yield 1, "\n\n".join(paras)

    @staticmethod
    def _extract_pptx(path: Path):
        from pptx import Presentation
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text_frame.text for shape in slide.shapes
                     if shape.has_text_frame and shape.text_frame.text.strip()]
            yield i, "\n\n".join(texts)
